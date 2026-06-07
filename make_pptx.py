#!/usr/bin/env python3
"""Build a HYBRID PowerPoint of the English deck.

Goal (per user): match the HTML look, but don't make text slides into pictures.

  * Slides that have a FIGURE/VIDEO (Introduction, Method, IK, Torch, Positioner,
    Results, + the designed title slide) -> a full-bleed screenshot of the
    rendered HTML slide, so they match the reveal.js deck EXACTLY. The Results
    slide additionally gets the real comparison video embedded (plays in PPT).
  * Slides that are TEXT-ONLY (System overview, Contribution, Questions) -> real,
    editable native PowerPoint text, styled CTU-blue / Lato to look like the HTML
    (no picture used where none is needed).

Screenshots come from headless Chrome against presentation_en.html (reveal.js
auto-scales each slide to fit 1280x720 -> one clean image per slide). Native
text + numbering are read from presentation_en.qmd. Czech speaker notes (from
the HTML) are attached to every slide.

Run via build.sh (after the HTML render). Needs: headless Chrome, python-pptx.
"""
import html
import os
import re
import shutil
import subprocess
import sys
import tempfile
from html.parser import HTMLParser

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

HERE = os.path.dirname(os.path.abspath(__file__))
QMD = os.path.join(HERE, "presentation_en.qmd")
HTML = os.path.join(HERE, "presentation_en.html")
OUT = os.path.join(HERE, "presentation_en.pptx")

RESULTS_VIDEO = os.path.join(HERE, "figures", "results_compare.mp4")
RESULTS_POSTER = os.path.join(HERE, "figures", "results_compare_poster.png")

SLIDE_W_PX, SLIDE_H_PX = 1280, 720
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5

# Colours matched to custom.scss.
CTU_BLUE = RGBColor(0x00, 0x65, 0xBD)
HEADING = RGBColor(0x1A, 0x1A, 0x1A)
BODY = RGBColor(0x22, 0x22, 0x22)
RULE = RGBColor(0xB8, 0xD0, 0xE8)   # light-blue h2 underline
FONT = "Lato"

CLEAN_CSS = """
<style id="pptx-screenshot-clean">
  .slide-menu-button { display: none !important; }
  .reveal .controls  { display: none !important; }
  .reveal .progress  { display: none !important; }
</style>
"""


def find_chrome():
    for c in ("google-chrome", "google-chrome-stable", "chromium",
              "chromium-browser", "brave-browser"):
        if shutil.which(c):
            return shutil.which(c)
    sys.exit("ERROR: no Chrome/Chromium found — cannot screenshot slides.")


# --------------------------------------------------------------------------
# Read the .qmd: per-content-slide (title, body-items). Title slide is YAML.
# --------------------------------------------------------------------------
def qmd_sections():
    src = open(QMD, encoding="utf-8").read()
    src = re.sub(r"<!--.*?-->", "", src, flags=re.S)      # drop commented slides
    parts = re.split(r"(?m)^##\s+(.*)$", src)
    secs = []
    for i in range(1, len(parts), 2):
        title = re.sub(r"\s*\{.*\}\s*$", "", parts[i]).strip()
        secs.append((title, parts[i + 1]))
    return secs


def _runs(text):
    """Split markdown text into (text, bold) runs on **bold** markers."""
    out = []
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if part:
            out.append((part, i % 2 == 1))
    return out


def parse_items(body):
    """Parse a slide body into (level, number_or_None, runs) list items."""
    body = body.split("::: {.notes}")[0]
    items, n0 = [], 0
    for raw in body.splitlines():
        s = raw.strip()
        if not s or s.startswith(":::"):
            continue
        indent = len(raw) - len(raw.lstrip())
        m0 = re.match(r"^(\d+)\.\s+(.*)$", s)
        m1 = re.match(r"^[-*]\s+(.*)$", s)
        if m0 and indent == 0:
            n0 += 1
            items.append((0, n0, _runs(m0.group(2))))
        elif m1:
            items.append((1, None, _runs(m1.group(1))))
    return items


# --- Czech speaker notes from the rendered HTML ---------------------------
class _Asides(HTMLParser):
    SKIP = {"style", "script", "svg", "mjx-container"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.notes, self._in, self._buf, self._skip = [], False, [], 0

    def handle_starttag(self, tag, attrs):
        if tag == "aside" and "notes" in (dict(attrs).get("class") or ""):
            self._in, self._buf, self._skip = True, [], 0
        elif self._in:
            if tag in self.SKIP:
                self._skip += 1
            elif self._skip == 0 and tag in ("br", "p", "div", "li"):
                self._buf.append("\n")

    def handle_endtag(self, tag):
        if tag == "aside" and self._in:
            self._in = False
            self.notes.append("".join(self._buf))
        elif self._in and tag in self.SKIP and self._skip:
            self._skip -= 1

    def handle_data(self, data):
        if self._in and self._skip == 0:
            self._buf.append(data)


def _clean(t):
    for pat in (r"<!--.*?-->", r"<!–.*?–>"):
        t = re.sub(pat, "", t, flags=re.S)
    lines = [ln.strip() for ln in t.split("\n")]
    out = []
    for ln in lines:
        if ln or (out and out[-1]):
            out.append(ln)
    return "\n".join(out).strip()


def deck_data():
    """Return (n_slides, notes[], is_figure[], sections[]) aligned by index."""
    raw = open(HTML, encoding="utf-8").read()
    m = re.search(r'data-notes="([^"]*)"', raw)
    title_note = _clean(html.unescape(m.group(1))) if m else ""
    a = _Asides()
    a.feed(raw)
    notes = [title_note] + [_clean(x) for x in a.notes]

    # figure vs text-only (+ heading text + uncounted flag), per slide section
    is_fig = []
    heads = []
    uncounted = []                                     # data-visibility="uncounted"
    for mm in re.finditer(r"(<section[^>]*>)(.*?)</section>", raw, re.S):
        tag, seg = mm.group(1), mm.group(2)
        if re.search(r"<h[12]", seg):                 # a real slide section
            is_fig.append(bool(re.search(r"<img |<video ", seg)))
            hm = re.search(r"<h[12][^>]*>(.*?)</h[12]>", seg, re.S)
            heads.append(re.sub(r"<[^>]+>", "", hm.group(1)).strip() if hm else "")
            uncounted.append('data-visibility="uncounted"' in tag)
    is_fig[0] = True                                   # title slide -> image
    uncounted[0] = False

    n = len(is_fig)
    notes = (notes + [""] * n)[:n]
    secs = qmd_sections()                              # content slides only
    return n, notes, is_fig, secs, heads, uncounted


# --------------------------------------------------------------------------
# Screenshots (only the slides we render as images)
# --------------------------------------------------------------------------
def screenshot(chrome, indices):
    tmp = tempfile.mkdtemp(prefix="deck_pptx_")
    src = open(HTML, encoding="utf-8").read().replace(
        "</head>", CLEAN_CSS + "</head>", 1)
    page = os.path.join(tmp, "deck.html")
    open(page, "w", encoding="utf-8").write(src)
    shots = {}
    for i in indices:
        out = os.path.join(tmp, f"slide-{i:02d}.png")
        subprocess.run([
            chrome, "--headless=new", "--disable-gpu", "--no-sandbox",
            "--hide-scrollbars", "--force-device-scale-factor=1",
            f"--window-size={SLIDE_W_PX},{SLIDE_H_PX}",
            "--virtual-time-budget=12000", f"--screenshot={out}",
            f"file://{page}#/{i}",
        ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        if not os.path.exists(out):
            sys.exit(f"ERROR: screenshot failed for slide {i}")
        shots[i] = out
        print(f"  shot slide {i + 1}")
    return tmp, shots


# --------------------------------------------------------------------------
# Native (editable) text slide, styled to look like the HTML
# --------------------------------------------------------------------------
def _set_marl(p, inches):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(inches))))
    pPr.set("indent", str(-int(Inches(0.28))))


def add_text_slide(prs, blank, title, items, page_no, total):
    slide = prs.slides.add_slide(blank)

    # Title (dark, bold, Lato) + a light-blue underline rule (matches h2).
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.30),
                                  Inches(12.2), Inches(0.95))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size, r.font.bold, r.font.name = Pt(34), True, FONT
    r.font.color.rgb = HEADING
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(0.55), Inches(1.28),
                                    Inches(12.78), Inches(1.28))
    ln.line.color.rgb = RULE
    ln.line.width = Pt(1.2)

    # Body
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55),
                                   Inches(12.0), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for level, num, runs in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7 if level == 0 else 3)
        _set_marl(p, 0.15 if level == 0 else 0.85)
        size = 19 if level == 0 else 15
        if num is not None:                            # "1. " prefix
            pr = p.add_run()
            pr.text = f"{num}. "
            pr.font.size, pr.font.name, pr.font.color.rgb = Pt(size), FONT, BODY
        for text, bold in runs:
            rn = p.add_run()
            rn.text = text
            rn.font.size, rn.font.name, rn.font.bold = Pt(size), FONT, bold
            rn.font.color.rgb = CTU_BLUE if bold else BODY

    if page_no is not None:
        _add_page_no(slide, page_no, total)
    return slide


def _add_page_no(slide, page_no, total):
    nb = slide.shapes.add_textbox(Inches(11.9), Inches(7.02),
                                  Inches(1.3), Inches(0.4))
    pr = nb.text_frame.paragraphs[0]
    pr.alignment = PP_ALIGN.RIGHT
    r = pr.add_run()
    r.text = f"{page_no} / {total}"
    r.font.size, r.font.name = Pt(11), FONT
    r.font.color.rgb = CTU_BLUE


def build():
    chrome = find_chrome()
    n, notes, is_fig, secs, heads, uncounted = deck_data()
    img_idx = [i for i in range(n) if is_fig[i]]
    n_counted = sum(1 for i in range(n) if not uncounted[i])   # total for "N / t"
    # Locate the Results slide by its heading (robust to extra slides such as
    # the split Questions deck), not by a fixed position.
    results_idx = next(
        (i for i, h in enumerate(heads) if h.strip().lower().startswith("results")),
        None)
    print(f"Deck: {n} slides — images: {[i+1 for i in img_idx]}, "
          f"text: {[i+1 for i in range(n) if not is_fig[i]]} "
          f"(Results video on #{results_idx + 1 if results_idx is not None else '—'}).")

    tmp, shots = screenshot(chrome, img_idx)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    for i in range(n):
        if is_fig[i]:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(shots[i], 0, 0,
                                     width=prs.slide_width,
                                     height=prs.slide_height)
            if i == results_idx and os.path.exists(RESULTS_VIDEO) \
                    and os.path.exists(RESULTS_POSTER):
                slide.shapes.add_movie(
                    RESULTS_VIDEO, Inches(1.00), Inches(1.35),
                    Inches(11.33), Inches(5.37),
                    poster_frame_image=RESULTS_POSTER, mime_type="video/mp4")
        else:
            title, body = secs[i - 1]               # content slide (title=YAML)
            # uncounted slides (Questions) show no page number; counted slides
            # use i+1 (all uncounted slides sit at the end, so i+1 == count).
            page_no = None if uncounted[i] else i + 1
            slide = add_text_slide(prs, blank, title, parse_items(body),
                                   page_no, n_counted)
        if notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]

    prs.save(OUT)
    shutil.rmtree(tmp, ignore_errors=True)
    print(f"Wrote {os.path.basename(OUT)}  ({n} slides, "
          f"{os.path.getsize(OUT) / 1e6:.1f} MB) — text slides are editable; "
          "figure slides match the HTML; Results video plays on click.")


if __name__ == "__main__":
    build()
